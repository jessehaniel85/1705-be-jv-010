#!/usr/bin/env python3
# Gera os slides .pptx das aulas BE-JV-010 usando o TEMPLATE OFICIAL Ada.
from pptx import Presentation
from pptx.oxml.ns import qn

TEMPLATE = "/Users/jessehaniel/Documents/adatech-educacional/Labs/assets/Ada _ Template PPT.pptx"
OUT = "/Users/jessehaniel/Documents/adatech-educacional/Labs/Caixa-EmbarqueTI-NovaProposta/planejamento-de-aulas/1705-be-jv-010/plano-de-aulas"

# Layouts do template oficial:
L_TITLE   = 0   # placeholders idx 0 (título) e 2 (subtítulo-título)
L_CONTENT = 2   # idx 0 (título), 1 (corpo), 2 (subtítulo)
L_POINT   = 1   # idx 0 (título único) — frase-síntese / divisor

def new_deck():
    prs = Presentation(TEMPLATE)
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)
    return prs

def _set_body(ph, lines):
    if isinstance(lines, str):
        lines = [lines]
    tf = ph.text_frame
    tf.clear()
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if ln.startswith("  - "):
            p.level = 1
            p.text = ln[4:]
        else:
            p.text = ln

def add(prs, layout_idx, fills):
    """fills: dict idx -> str (texto) | list (bullets)."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    by_idx = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    for idx, val in fills.items():
        ph = by_idx.get(idx)
        if ph is None:
            continue
        if isinstance(val, list):
            _set_body(ph, val)
        else:
            ph.text = val
    return slide

def titulo(t, sub):       return (L_TITLE,   {0: t, 2: sub})
def conteudo(t, body, sub=None):
    f = {0: t, 1: body}
    if sub: f[2] = sub
    return (L_CONTENT, f)
def ponto(t):             return (L_POINT,   {0: t})

def build(filename, slides):
    prs = new_deck()
    for layout_idx, fills in slides:
        add(prs, layout_idx, fills)
    path = f"{OUT}/{filename}"
    prs.save(path)
    print("ok:", filename, "slides:", len(slides))

# ============================ AULA 1 — Apresentação do Curso ============================
aula1 = [
 titulo("Arquitetura de Software e Ágil II", "BE-JV-010 · EmbarqueTI · Nível III · Apresentação do Curso"),
 conteudo("Agenda", ["Metodologia Ada", "Conteúdo do módulo", "Cronograma (9 encontros)",
    "Os dois projetos", "Avaliação e rubrica", "IA & agentes: por que entram aqui", "Combinados"]),
 conteudo("Metodologia Ada", [
    "Aprendizagem ativa: PBL + Sala de Aula Invertida (adaptada)",
    "1ª metade: problema → discussão → solução ao vivo",
    "2ª metade: você evolui a solução + trabalha no projeto em grupo",
    "  - Estudo e projeto acontecem DENTRO da aula",
    "Turma sênior: debate de alto nível + ponte com o legado"]),
 conteudo("Conteúdo do módulo", [
    "Microsserviços e DDD (fronteiras de domínio)",
    "Consistência distribuída e padrão SAGA",
    "Cache compartilhado (Redis)",
    "Mensageria: filas (RabbitMQ) e tópicos (Kafka)",
    "Resiliência: retry, circuit breaker, DLQ",
    "Contract testing (PACT)",
    "  - Fio transversal: IA & agentes em cada tema"]),
 conteudo("Cronograma", [
    "22/06 Abertura + DDD e fronteiras",
    "24/06 Consistência distribuída e SAGA",
    "26/06 Cache com Redis",
    "29/06 Comunicação assíncrona (producer/consumer)",
    "01/07 Filas com RabbitMQ + DLQ",
    "03/07 Tópicos/eventos com Kafka",
    "06/07 Resiliência e retry",
    "08/07 Contract testing (PACT)",
    "10/07 Bancas dos projetos + devolutiva"]),
 conteudo("Projeto guiado (em aula)", [
    "API de Comprovantes PIX",
    "Construída ao vivo, incremento a cada aula",
    "POST → 202 → fila → gravação; consulta com cache + fallback",
    "  - É a referência dos padrões — não é avaliado"]),
 conteudo("Projeto final (em grupo)", [
    "Tema de LIVRE escolha do grupo",
    "O que amarra são os CRITÉRIOS de avaliação, não o tema",
    "Grupos de 3 a 5 (sugestão); escopo escala com o tamanho",
    "Desenvolvido na 2ª metade das aulas; defesa na banca (10/07)",
    "  - Formem os grupos até a próxima aula"]),
 conteudo("Avaliação e rubrica", [
    "Rubrica de 9 critérios (decomposição, async, cache, resiliência, testes, decisões...)",
    "Avalia DOMÍNIO do padrão e DECISÃO arquitetural — não a infra disponível",
    "Banca de defesa + auto-avaliação",
    "Listas de exercícios e participação"]),
 conteudo("Avalie a satisfação", [
    "Toda aula é possível (e muito importante!) avaliar no LMS",
    "Seu NPS orienta os ajustes da próxima aula",
    "Meta de excelência Ada: 90+"]),
 conteudo("IA & agentes: por que entram aqui", [
    "O módulo é de 2024; a realidade de quem opera esses sistemas mudou",
    "Cada padrão clássico tem um paralelo direto com IA/agentes",
    "  - Bounded context = fronteira de agente; SAGA = agentic workflow",
    "  - Cache = semantic caching/RAG; tópicos = memória de agente",
    "Não é hype: sempre ancorado no padrão de engenharia"]),
 conteudo("Combinados", [
    "Início, intervalo e término (acordar horários da turma)",
    "Estar online ≠ estar presente — câmera e participação contam",
    "Ambiente: tudo via Nexus corporativo; sem depender de Docker/sites externos",
    "Dúvidas são bem-vindas durante o live coding"]),
 # ---- Conteúdo de DDD (âncoras visuais; o aprofundamento está no material do aluno) ----
 ponto("DDD: a ferramenta para decidir ONDE cortar o sistema"),
 conteudo("O problema do corte", [
    "Cortar no lugar errado = monólito distribuído (o pior dos dois mundos)",
    "A pergunta certa não é 'em quantos serviços eu divido?'",
    "É 'quais fronteiras mudam por razões e ritmos diferentes?'",
    "Decisão de DOMÍNIO, não de tecnologia"]),
 conteudo("DDD em dois níveis", [
    "Estratégico — ONDE estão as fronteiras (viram microsserviços)",
    "  - domínio, subdomínios, linguagem ubíqua, bounded context, event storming",
    "Tático — COMO modelar dentro de cada fronteira",
    "  - entidades, value objects, agregados",
    "Não é arquitetura em camadas; não é framework"]),
 conteudo("Blocos táticos", [
    "Entidade — tem identidade e ciclo de vida (ex.: Comprovante, por id)",
    "Value Object — imutável, igualdade por valor (ex.: ChavePix, Dinheiro)",
    "Agregado + raiz — cluster que muda junto e guarda a invariante",
    "  - 1 transação = 1 agregado; referência entre agregados por id",
    "Regra de ouro: o agregado NUNCA atravessa fronteira de serviço"]),
 conteudo("Bounded Context & Context Map", [
    "Mesmo termo, significados diferentes por contexto (o 'Cliente' de cada área)",
    "Bounded context = fronteira do modelo e da linguagem ubíqua",
    "1 contexto → 1 ou mais microsserviços (nunca o contrário)",
    "Context map: customer/supplier, conformist, ACL, OHS/published language"]),
 conteudo("Event Storming", [
    "Oficina com post-its: negócio + técnicos modelam o fluxo juntos",
    "🟧 evento · 🟦 comando · 🟨 agregado · 🟪 política · 🧍 ator · 🟥 hotspot",
    "Eventos → comandos/atores → agregados → políticas → contextos",
    "  - Políticas viram mensageria; hotspots viram requisitos de resiliência"]),
 conteudo("Dois cenários (no material do aluno)", [
    "1) Comprovantes PIX — do relato do negócio aos microsserviços, ao vivo",
    "2) Fatura de cartão — invariante 'soma dos lançamentos = total'",
    "Mesmo método; consistência FORTE dentro do agregado, EVENTUAL entre contextos"]),
 ponto("Microsserviço não é tamanho de código — é fronteira de domínio com dono."),
]

# ============================ AULAS 2–8 (decks enxutos) ============================
def lean(num, d):
    return [
      titulo(d["titulo"], f"BE-JV-010 · Aula {num} · {d['data']} · {d['competencia']}"),
      conteudo("Problema gerador", d["problema"]),
      conteudo("Conceitos-chave", d["conceitos"]),
      conteudo("Ponte com o legado", d["legado"]),
      conteudo("Ângulo IA & agentes", d["ia"]),
      conteudo("Desafio (studio)", d["desafio"]),
      ponto(d["sintese"]),
    ]

aulas = {
 2: {"titulo":"Consistência distribuída e SAGA","data":"24/06","competencia":"Padrão SAGA",
     "problema":["A emissão respondeu 202, mas a gravação falhou.",
                 "O cliente tem um comprovante que não existe no banco.","Como evitar — e como compensar?"],
     "conceitos":["Por que o 2PC saiu de cena","Consistência eventual e seu custo ao negócio",
                  "SAGA: orquestração × coreografia","Ações compensatórias","Outbox pattern; idempotência"],
     "legado":["A 'orquestração' via stored procedure / job control do mainframe",
               "já é uma SAGA implícita — só que sem compensação explícita"],
     "ia":["Agentic workflow = SAGA com passos compensatórios",
           "Orquestrador central × agentes reagindo a eventos","Reexecutar passo não-determinístico exige idempotência"],
     "desafio":["Base: compensação + idempotência por chave (sem duplicar)",
                "Stretch: versão coreografada por eventos; onde entra o outbox?"],
     "sintese":"Em distribuído você não escolhe ter falha parcial — só se vai tratá-la de propósito."},
 3: {"titulo":"Cache com Redis","data":"26/06","competencia":"Cache compartilhado",
     "problema":["A consulta de comprovante é muito mais lida que escrita","e está martelando o banco.",
                 "Como aliviar sem servir dado inconsistente?"],
     "conceitos":["Cache local × distribuído (por que Redis e não um HashMap)",
                  "cache-aside, write-through, write-behind","TTL e invalidação","Stampede / thundering herd"],
     "legado":["Cache em tabela de banco e sessão pegajosa no app server",
               "Redis tira o estado do nó e resolve replicação de sessão"],
     "ia":["Semantic caching: cachear resposta de LLM por similaridade de embedding",
           "Redis como vector store para RAG","Cache deixa de ser otimização e vira requisito econômico"],
     "desafio":["Base: TTL + medir hit/miss",
                "Stretch: invalidação ao atualizar + mitigar stampede (lock/jitter)"],
     "sintese":"Cache não é guardar tudo — é decidir o que pode estar um pouco velho."},
 4: {"titulo":"Comunicação assíncrona","data":"29/06","competencia":"Producer/consumer",
     "problema":["A gravação é lenta e às vezes o banco está fora.",
                 "O cliente não pode esperar nem perder o comprovante.","Como aceitar agora e processar depois?"],
     "conceitos":["Mensagem × evento","at-most/at-least/exactly-once (e por que exactly-once é mito)",
                  "Acoplamento temporal","Idempotência do consumidor"],
     "legado":["Quem usou IBM MQ / JMS / MQSeries já fez producer/consumer",
               "Muda o ferramental, não o conceito"],
     "ia":["Desacoplar inferência cara do caminho síncrono","Filas de tarefas de agentes (workers)",
           "Human-in-the-loop assíncrono"],
     "desafio":["Base: consumidor idempotente (reprocesso não duplica)",
                "Stretch: simular consumidor lento/caído; discutir backpressure"],
     "sintese":"Assíncrono não é 'mais rápido' — é 'não falhar junto'."},
 5: {"titulo":"Filas com RabbitMQ + DLQ","data":"01/07","competencia":"Filas",
     "problema":["A gravação falha de forma intermitente (banco oscila)",
                 "e às vezes permanente (mensagem malformada).","Os dois casos não podem ser tratados igual."],
     "conceitos":["Exchange, binding, routing key, queue","Tipos de exchange",
                  "Retry × Dead Letter Queue","Ack/nack e redelivery no crash"],
     "legado":["DLQ é a 'fila de exceção / arquivo de rejeitados' do batch",
               "agora explícita e observável"],
     "ia":["Work queue para jobs de IA respeitando rate limit","Backpressure quando o LLM é o gargalo",
           "DLQ para chamadas de IA que falham (sem perda silenciosa)"],
     "desafio":["Base: configurar DLQ e provocar mensagem envenenada",
                "Stretch: retry com backoff + idempotência sob redelivery"],
     "sintese":"Mensagem nunca some por acidente — ou foi processada, ou está numa DLQ por decisão."},
 6: {"titulo":"Tópicos / eventos com Kafka","data":"03/07","competencia":"Tópicos / pub-sub",
     "problema":["Depois de gravar, 3 áreas querem reagir: notificação, antifraude e BI.",
                 "O gravador não pode conhecer todas.","Como avisar 'aconteceu' sem acoplar?"],
     "conceitos":["Fila × tópico","Log de eventos, offset, consumer group, partição",
                  "Event sourcing e replay","Ordenação por chave"],
     "legado":["O 'arquivo de movimento do dia' lido por vários jobs noturnos",
               "já era um log de eventos batch — Kafka é isso em streaming"],
     "ia":["Arquitetura event-driven para agentes","Event sourcing como memória reproduzível do agente",
           "Streaming para ingestão/RAG contínuos"],
     "desafio":["Base: plugar um 3º consumer group que lê desde o início",
                "Stretch: particionamento/ordenação por chave e o que quebra"],
     "sintese":"Fila entrega trabalho; tópico publica história."},
 7: {"titulo":"Resiliência: retry e circuit breaker","data":"06/07","competencia":"Resiliência e retry",
     "problema":["A notificação depende de um gateway externo que oscila.",
                 "Queremos reentregar sem inundar","e parar de tentar quando claramente está fora."],
     "conceitos":["Backoff exponencial + jitter (evitar retry storm)","Idempotência como pré-condição",
                  "Circuit breaker: closed/open/half-open","Timeout e bulkhead","@RetryableTopic + DLT"],
     "legado":["'Reprocessamento do batch que deu erro' e limites de retentativa do scheduler",
               "agora reativos e por requisição"],
     "ia":["APIs de LLM são instáveis por natureza (rate limit, timeout, 5xx)",
           "Fallback de modelo quando o circuito abre","Idempotência com saída não-determinística"],
     "desafio":["Base: @RetryableTopic com backoff até a DLT",
                "Stretch: retry + breaker + idempotência sob falha intermitente E permanente"],
     "sintese":"Resiliência é decidir como falhar de propósito, antes que o sistema decida por você."},
 8: {"titulo":"Contract testing (PACT)","data":"08/07","competencia":"Arquitetura testável",
     "problema":["Times diferentes mantêm emissor e gravador.",
                 "Uma mudança de contrato quebra o outro em produção.","Como travar isso no CI?"],
     "conceitos":["Pirâmide de testes; por que E2E não escala","Consumer-driven contracts",
                  "PACT: consumidor gera, provedor verifica","can-I-deploy; arquitetura amigável a testes"],
     "legado":["O 'contrato' que vivia num documento Word / layout de arquivo",
               "agora é executável e roda no pipeline"],
     "ia":["Testar saída não-determinística por schema/propriedade","Evals como contrato de um componente de IA",
           "Structured outputs / JSON schema como contrato"],
     "desafio":["Base: contract test que falha ao mudar o contrato",
                "Stretch: verificação de schema/propriedade (estilo eval) no mvn verify"],
     "sintese":"Contrato bom é o que falha no seu CI, não na produção do outro time."},
}

# ============================ AULA 9 — Bancas ============================
aula9 = [
 titulo("Bancas dos projetos + Devolutiva", "BE-JV-010 · Aula 9 · 10/07 · Defesa de arquitetura (Nível III)"),
 conteudo("Como funciona a banca", [
    "Demo do que roda (declare o perfil A/B/C e os fallbacks)",
    "Defenda as DECISÕES, não só o que o sistema faz",
    "Arguição do docente + perguntas dos pares",
    "Reflexão: como o grupo usou IA e o que validou à mão"]),
 conteudo("O que a banca avalia", [
    "Decomposição de domínio e bases segregadas",
    "Assíncrono, idempotência e consistência",
    "Cache, resiliência e contract testing",
    "Qualidade das decisões (ADRs)",
    "  - Não penaliza restrição de ambiente"]),
 conteudo("Devolutiva e auto-avaliação", [
    "Auto-avaliação por competência do módulo",
    "Padrões fortes e fracos observados nas bancas",
    "Entrega final do repositório (com AVALIACAO.md)"]),
 ponto("Sistemas distribuídos: transformar falha inevitável em comportamento previsível."),
 conteudo("IA e o papel do arquiteto", [
    "O que a IA acelera — e o que ela não decide",
    "Por que domínio de arquitetura distribuída fica MAIS valioso na era dos agentes",
    "Avalie a satisfação no LMS — obrigado pela jornada!"]),
]

build("aula-01-slides.pptx", aula1)
for n, d in aulas.items():
    build(f"aula-0{n}-slides.pptx", lean(n, d))
build("aula-09-slides.pptx", aula9)
print("CONCLUÍDO")
